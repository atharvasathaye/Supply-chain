from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Supply Chain Pressure vs. Inflation"
    subtitle.text = "A Predictive Machine Learning Analysis\n\nBy Atharva Sathaye"

    # Slide 1: Introduction
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Introduction"
    tf = body_shape.text_frame
    tf.text = "The modern global economy is highly interconnected."
    p = tf.add_paragraph()
    p.text = "Goal: Investigate if supply chain bottlenecks act as a leading indicator for consumer inflation."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Data Sources: Federal Reserve (CPI) and NY Fed (GSCPI)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Deliverables: End-to-end data pipeline, statistical proofs, and an interactive ML dashboard."
    p.level = 1

    # Slide 2: Background & Analogy
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Background & The 'Traffic Jam' Analogy"
    tf = body_shape.text_frame
    tf.text = "Analogy: A traffic jam on a highway."
    p = tf.add_paragraph()
    p.text = "When an accident occurs (supply chain shock), traffic doesn't stop immediately miles back."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "It takes time for the ripple effect to hit the cars behind (the consumers)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "In economics, a shipping delay in China takes 1-6 months to cause price hikes in US grocery stores."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "This delay creates a window where we can predict future inflation before it happens."
    p.level = 1

    # Slide 3: Methodology (Flowchart substitute)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Methodology & Architecture"
    tf = body_shape.text_frame
    tf.text = "1. Automated Data Engineering Pipeline"
    p = tf.add_paragraph()
    p.text = "Python script fetches live data via APIs and normalizes dates."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Statistical Proof (Granger Causality)"
    p = tf.add_paragraph()
    p.text = "Statistically proved that GSCPI changes preceded CPI changes (p-values < 0.05)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Machine Learning Forecasting (VAR Model)"
    p = tf.add_paragraph()
    p.text = "Vector Autoregression predicts inflation 6 months into the future."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Cloud Deployment (MLOps)"
    p = tf.add_paragraph()
    p.text = "Dockerized Streamlit app with GitHub Actions for automated monthly updates."
    p.level = 1

    # Slide 4: Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Conclusion & Impact"
    tf = body_shape.text_frame
    tf.text = "Supply chains are mathematical predictors of inflation."
    p = tf.add_paragraph()
    p.text = "Physical goods (Vehicles, Food) react violently to supply shocks, while services remain stable."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Business Impact:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Companies can use these models to adjust pricing strategies months before inflation hits their P&L."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Demonstrates full-stack data capabilities from engineering to machine learning to BI dashboarding."
    p.level = 1

    # Save
    base_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(base_dir, "Project_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == '__main__':
    create_presentation()
